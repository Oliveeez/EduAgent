#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
调试脚本：检查布局重叠问题
用于诊断bbox提取和VLM优化是否正确工作
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# 添加项目根目录到 Python 路径
project_root = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(project_root))


def check_bbox_accuracy():
    """检查1：bbox提取是否准确"""
    print("=" * 80)
    print("🔍 检查1：BBox提取准确性")
    print("=" * 80)
    print()
    
    base_dir = Path(project_root) / "data" / "pipeline_outputs" / "test_ppt_only"
    pptx_path = base_dir / "pptx" / "presentation.pptx"
    bboxes_path = base_dir / "bboxes.json"
    
    # 1. 读取实际PPT
    print("📂 读取PPT文件...")
    prs = Presentation(str(pptx_path))
    
    # 2. 读取提取的bboxes
    with open(bboxes_path, 'r', encoding='utf-8') as f:
        bboxes_data = json.load(f)
    
    print(f"   ✅ 共 {len(prs.slides)} 页PPT")
    print()
    
    # 3. 逐页对比
    for slide_idx, slide in enumerate(prs.slides):
        bbox_slide = None
        for bs in bboxes_data.get('slides', []):
            if bs.get('slide_id') == slide_idx:
                bbox_slide = bs
                break
        
        if not bbox_slide:
            continue
        
        print(f"📄 Slide {slide_idx}:")
        print(f"   PPT实际元素数: {len(slide.shapes)}")
        print(f"   BBox提取数量: {len(bbox_slide.get('elements', []))}")
        
        # 对比每个元素
        ppt_shapes = [s for s in slide.shapes if not s.is_placeholder]
        bbox_elements = bbox_slide.get('elements', [])
        
        if len(ppt_shapes) != len(bbox_elements):
            print(f"   ⚠️ 数量不匹配！")
        
        for i, shape in enumerate(ppt_shapes):
            actual_left = shape.left / 914400  # EMU to inches
            actual_top = shape.top / 914400
            actual_width = shape.width / 914400
            actual_height = shape.height / 914400
            
            if i < len(bbox_elements):
                bbox = bbox_elements[i]
                bbox_left = bbox.get('left')
                bbox_top = bbox.get('top')
                bbox_width = bbox.get('width')
                bbox_height = bbox.get('height')
                
                # 计算误差
                error_left = abs(actual_left - bbox_left) if bbox_left else 999
                error_top = abs(actual_top - bbox_top) if bbox_top else 999
                error_width = abs(actual_width - bbox_width) if bbox_width else 999
                error_height = abs(actual_height - bbox_height) if bbox_height else 999
                
                max_error = max(error_left, error_top, error_width, error_height)
                
                print(f"   元素 {i} ({bbox.get('type', '未知')}):")
                print(f"      PPT shape: L={actual_left:.2f}, T={actual_top:.2f}, W={actual_width:.2f}, H={actual_height:.2f}")
                print(f"      BBox提取: L={bbox_left:.2f}, T={bbox_top:.2f}, W={bbox_width:.2f}, H={bbox_height:.2f}")
                
                # 对于文本元素，检查高度增加（表示使用了实际渲染高度）
                if bbox.get('type') == 'text' and bbox_height > actual_height + 0.05:
                    height_increase = bbox_height - actual_height
                    print(f"      📏 文本高度调整: +{height_increase:.2f} 英寸 (考虑实际渲染高度)")
                
                if max_error > 0.1:  # 误差超过0.1英寸
                    print(f"      ❌ 位置误差过大: {max_error:.2f} 英寸")
                else:
                    print(f"      ✅ 位置误差可接受: {max_error:.3f} 英寸")
        
        print()


def check_overlaps():
    """检查2：检测重叠"""
    print("=" * 80)
    print("🔍 检查2：检测元素重叠")
    print("=" * 80)
    print()
    
    base_dir = Path(project_root) / "data" / "pipeline_outputs" / "test_ppt_only"
    
    # 检查原始PPT
    print("📂 检查原始PPT (presentation.pptx)...")
    check_pptx_overlaps(base_dir / "pptx" / "presentation.pptx")
    print()
    
    # 检查优化后PPT
    print("📂 检查优化后PPT (presentation_optimized.pptx)...")
    check_pptx_overlaps(base_dir / "pptx" / "presentation_optimized.pptx")
    print()


def check_pptx_overlaps(pptx_path: Path):
    """检查单个PPTX的重叠情况"""
    prs = Presentation(str(pptx_path))
    
    total_overlaps = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        shapes = [s for s in slide.shapes if not s.is_placeholder]
        
        overlaps = []
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                if shapes_overlap(shapes[i], shapes[j]):
                    overlaps.append((i, j))
        
        if overlaps:
            print(f"   ⚠️ Slide {slide_idx}: 发现 {len(overlaps)} 处重叠")
            for i, j in overlaps:
                s1 = shapes[i]
                s2 = shapes[j]
                print(f"      元素{i} ({get_shape_info(s1)}) 与 元素{j} ({get_shape_info(s2)}) 重叠")
            total_overlaps += len(overlaps)
    
    if total_overlaps == 0:
        print(f"   ✅ 未发现重叠")
    else:
        print(f"   总计: {total_overlaps} 处重叠")


def shapes_overlap(s1, s2) -> bool:
    """判断两个shape是否重叠"""
    left1 = s1.left / 914400
    top1 = s1.top / 914400
    right1 = (s1.left + s1.width) / 914400
    bottom1 = (s1.top + s1.height) / 914400
    
    left2 = s2.left / 914400
    top2 = s2.top / 914400
    right2 = (s2.left + s2.width) / 914400
    bottom2 = (s2.top + s2.height) / 914400
    
    # 检测是否重叠
    return not (right1 <= left2 or right2 <= left1 or bottom1 <= top2 or bottom2 <= top1)


def get_shape_info(shape) -> str:
    """获取shape的简要信息"""
    left = shape.left / 914400
    top = shape.top / 914400
    width = shape.width / 914400
    height = shape.height / 914400
    return f"L={left:.1f},T={top:.1f},W={width:.1f}x{height:.1f}"


def check_vlm_adjustments():
    """检查3：VLM调整是否合理"""
    print("=" * 80)
    print("🔍 检查3：VLM调整方案")
    print("=" * 80)
    print()
    
    base_dir = Path(project_root) / "data" / "pipeline_outputs" / "test_ppt_only"
    adjustments_path = base_dir / "layout_adjustments.json"
    
    if not adjustments_path.exists():
        print("   ❌ 未找到调整记录文件")
        return
    
    with open(adjustments_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    adjustments = data.get('adjustments', [])
    
    if not adjustments:
        print("   ⚠️ 没有任何调整记录！")
        print("   可能原因：")
        print("      1. VLM认为布局没问题（但实际有问题）")
        print("      2. VLM返回了建议但未被应用")
        print("      3. VLM调用失败")
        return
    
    print(f"   ✅ 共 {len(adjustments)} 条调整记录")
    print()
    
    # 按slide分组
    by_slide = {}
    for adj in adjustments:
        slide_idx = adj.get('slide_idx', -1)
        if slide_idx not in by_slide:
            by_slide[slide_idx] = []
        by_slide[slide_idx].append(adj)
    
    for slide_idx, slide_adjs in sorted(by_slide.items()):
        print(f"   📄 Slide {slide_idx}: {len(slide_adjs)} 个调整")
        for adj in slide_adjs:
            elem_idx = adj.get('element_idx')
            action = adj.get('action')
            reason = adj.get('reason', '')
            
            if action == 'move':
                new_left = adj.get('new_left')
                new_top = adj.get('new_top')
                print(f"      元素{elem_idx}: 移动到 ({new_left:.2f}, {new_top:.2f})")
            elif action == 'resize':
                new_width = adj.get('new_width')
                new_height = adj.get('new_height')
                print(f"      元素{elem_idx}: 调整大小为 {new_width:.2f}x{new_height:.2f}")
            
            if reason:
                print(f"         原因: {reason[:60]}...")
        print()


def check_out_of_bounds():
    """检查4：检查超出边界的元素"""
    print("=" * 80)
    print("🔍 检查4：检查超出边界")
    print("=" * 80)
    print()
    
    PAGE_WIDTH = 13.33
    PAGE_HEIGHT = 7.5
    SAFE_MARGIN = 0.5
    
    base_dir = Path(project_root) / "data" / "pipeline_outputs" / "test_ppt_only"
    
    # 检查原始PPT
    print(f"📂 检查原始PPT...")
    check_pptx_bounds(base_dir / "pptx" / "presentation.pptx", PAGE_WIDTH, PAGE_HEIGHT, SAFE_MARGIN)
    print()
    
    # 检查优化后PPT
    print(f"📂 检查优化后PPT...")
    check_pptx_bounds(base_dir / "pptx" / "presentation_optimized.pptx", PAGE_WIDTH, PAGE_HEIGHT, SAFE_MARGIN)
    print()


def check_pptx_bounds(pptx_path: Path, page_width: float, page_height: float, safe_margin: float):
    """检查PPTX中是否有超出边界的元素"""
    prs = Presentation(str(pptx_path))
    
    total_violations = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        shapes = [s for s in slide.shapes if not s.is_placeholder]
        
        violations = []
        for i, shape in enumerate(shapes):
            left = shape.left / 914400
            top = shape.top / 914400
            right = (shape.left + shape.width) / 914400
            bottom = (shape.top + shape.height) / 914400
            
            issues = []
            if left < safe_margin:
                issues.append(f"左边距{left:.2f}<{safe_margin}")
            if top < safe_margin:
                issues.append(f"上边距{top:.2f}<{safe_margin}")
            if right > page_width - safe_margin:
                issues.append(f"右边界{right:.2f}>{page_width-safe_margin}")
            if bottom > page_height:
                issues.append(f"底部{bottom:.2f}>{page_height}")
            
            if issues:
                violations.append((i, issues))
        
        if violations:
            print(f"   ⚠️ Slide {slide_idx}: {len(violations)} 个元素超出边界")
            for elem_idx, issues in violations:
                print(f"      元素{elem_idx}: {', '.join(issues)}")
            total_violations += len(violations)
    
    if total_violations == 0:
        print(f"   ✅ 所有元素都在安全区域内")
    else:
        print(f"   总计: {total_violations} 个元素超出边界")


def main():
    """主函数"""
    print()
    print("🔧 布局调试工具")
    print()
    
    # 检查1：bbox提取准确性
    check_bbox_accuracy()
    
    # 检查2：重叠检测
    check_overlaps()
    
    # 检查3：VLM调整方案
    check_vlm_adjustments()
    
    # 检查4：边界检查
    check_out_of_bounds()
    
    print("=" * 80)
    print("✅ 调试完成")
    print("=" * 80)
    print()
    print("💡 诊断建议:")
    print("   1. 如果'检查1'显示bbox提取不准确 → 问题在Stage 4 (bbox_extractor)")
    print("   2. 如果'检查2'原始PPT无重叠，优化后有重叠 → 问题在VLM调整逻辑")
    print("   3. 如果'检查3'没有调整记录 → 检查VLM是否正常工作")
    print("   4. 如果'检查3'有调整记录但'检查2'仍有重叠 → VLM建议不合理或应用失败")
    print()


if __name__ == "__main__":
    main()

