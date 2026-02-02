# stage5_vlm_optimizer.py
# Stage 5: VLM布局优化

import json
import re
from typing import List, Dict, Any, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

import sys
sys.path.append(str(Path(__file__).parent.parent.parent))

from utils.llm import CustomLLM
from .models import SlideBBoxes, BoundingBox


class VLMLayoutOptimizer:
    """
    VLM布局优化器
    
    功能：
    1. 将bounding box信息发送给LLM
    2. 获取布局调整建议
    3. 自动应用调整
    
    使用CustomLLM进行布局分析
    """
    
    # 页面尺寸
    PAGE_WIDTH = 13.33
    PAGE_HEIGHT = 7.5
    
    def __init__(self, pptx_path: str, output_dir: str):
        """
        初始化优化器
        
        Args:
            pptx_path: PPTX文件路径
            output_dir: 输出目录
        """
        self.pptx_path = Path(pptx_path)
        self.output_dir = Path(output_dir)
        self.prs = Presentation(str(self.pptx_path))
        self.llm = CustomLLM()
        
        self.adjustment_log: List[Dict] = []
    
    def optimize(self, bboxes: List[SlideBBoxes], max_iterations: int = 2, slides_data: Optional[List] = None) -> Path:
        """
        优化PPTX布局（优化版）
        
        优化策略：
        1. 智能跳过：简单页面（只有标题+文本）跳过VLM
        2. 提前终止：第1次迭代无问题则跳过第2次
        3. 简化检测：只检测有内容的元素重叠
        4. 并行调用：多页同时调用LLM
        
        Args:
            bboxes: 每页的bounding boxes（已过滤空元素）
            max_iterations: 最大迭代次数（默认2）
            slides_data: 幻灯片数据（用于判断是否需要优化）
        
        Returns:
            优化后的PPTX路径
        """
        print("  🔄 开始VLM布局优化...")
        
        # 筛选需要优化的页面（智能跳过）
        pages_to_optimize = []
        for slide_idx, slide_bboxes in enumerate(bboxes):
            # 获取对应的slide_data
            slide_data = slides_data[slide_idx] if slides_data and slide_idx < len(slides_data) else None
            
            if self._should_skip_page(slide_bboxes, slide_data):
                print(f"    ⏭️ 跳过页面 {slide_idx+1}（简单布局，无需优化）")
                continue
            pages_to_optimize.append((slide_idx, slide_bboxes))
        
        if not pages_to_optimize:
            print("    ✅ 所有页面都是简单布局，无需优化")
            output_path = self.output_dir / "pptx" / "presentation_optimized.pptx"
            output_path.parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(str(output_path))
            return output_path
        
        print(f"    📊 需要优化 {len(pages_to_optimize)} 页")
        
        # 迭代优化
        for iteration in range(max_iterations):
            print(f"    迭代 {iteration + 1}/{max_iterations}")
            
            # 并行处理多页
            needs_adjustment = self._optimize_pages_parallel(pages_to_optimize)
            
            # 提前终止：如果第1次迭代没发现问题，不进行第2次
            if not needs_adjustment:
                print("    ✅ 布局检查通过，提前终止")
                break
        
        # 保存优化后的PPTX
        output_path = self.output_dir / "pptx" / "presentation_optimized.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        
        # 保存调整日志
        self._save_adjustment_log()
        
        print(f"  ✅ 优化完成: {output_path}")
        return output_path
    
    def _should_skip_page(self, slide_bboxes: SlideBBoxes, slide_data=None) -> bool:
        """
        判断页面是否应该跳过VLM优化
        
        需要优化的情况：
        1. 有多个blocks（>2个）
        2. 有manim_relation、image、code等需要布局的元素
        3. 有重叠
        
        Args:
            slide_bboxes: 单页的bounding boxes
            slide_data: 单页的幻灯片数据（包含blocks信息）
        
        Returns:
            True表示应该跳过
        """
        # 如果有slide_data，优先基于blocks判断
        if slide_data and hasattr(slide_data, 'blocks'):
            blocks = slide_data.blocks
            
            # 如果只有1个block（通常是标题），可以跳过
            if len(blocks) <= 1:
                return True
            
            # 检查是否有需要布局优化的元素类型
            needs_optimization = False
            for block in blocks:
                block_type = block.block_type.value if hasattr(block.block_type, 'value') else str(block.block_type)
                if block_type in ["manim_relation", "image", "code", "formula", "manim_visual"]:
                    needs_optimization = True
                    break
            
            # 如果有需要优化的元素，或者blocks数量 > 2，则需要优化
            if needs_optimization or len(blocks) > 2:
                return False
        
        # 回退到基于bbox的判断
        has_visual_element = False
        text_count = 0
        
        for elem in slide_bboxes.elements:
            if elem.element_type == "image":
                has_visual_element = True
            elif elem.element_type == "text":
                text_count += 1
        
        # 如果有图片/GIF等视觉元素，需要优化
        if has_visual_element:
            return False
        
        # 如果只有1-2个文本框（标题+内容），可以跳过
        if text_count <= 2:
            return True
        
        # 如果有重叠，不能跳过
        if slide_bboxes.has_overlaps():
            return False
        
        return True
    
    def _optimize_pages_parallel(self, pages_to_optimize: List[tuple]) -> bool:
        """
        并行优化多页
        
        Args:
            pages_to_optimize: [(slide_idx, slide_bboxes), ...]
        
        Returns:
            是否有调整
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        needs_adjustment = False
        
        # 使用线程池并行处理（限制并发数避免API限流）
        with ThreadPoolExecutor(max_workers=3) as executor:
            futures = {
                executor.submit(self._optimize_single_page, slide_idx, slide_bboxes): slide_idx
                for slide_idx, slide_bboxes in pages_to_optimize
            }
            
            for future in as_completed(futures):
                slide_idx = futures[future]
                try:
                    had_adjustment = future.result()
                    if had_adjustment:
                        needs_adjustment = True
                except Exception as e:
                    print(f"      ⚠️ 页面 {slide_idx+1} 优化失败: {e}")
        
        return needs_adjustment
    
    def _optimize_single_page(self, slide_idx: int, slide_bboxes: SlideBBoxes) -> bool:
        """
        优化单页
        
        Args:
            slide_idx: 页面索引
            slide_bboxes: 该页的bounding boxes
        
        Returns:
            是否有调整
        """
        # 简化检测：只检查有内容的元素重叠
        if not slide_bboxes.has_overlaps() and not self._check_visual_issues(slide_bboxes):
            return False
        
        # 获取LLM建议
        adjustments = self._get_llm_suggestions(slide_bboxes)
        
        if adjustments:
            # 应用调整（传递slide_bboxes以获取shape_index）
            self._apply_adjustments(slide_idx, adjustments, slide_bboxes)
            return True
        
        return False
    
    def _check_visual_issues(self, slide_bboxes: SlideBBoxes) -> bool:
        """
        检查视觉问题
        
        Args:
            slide_bboxes: 单页的bounding boxes
        
        Returns:
            是否有问题
        """
        for elem in slide_bboxes.elements:
            # 检查是否超出边界
            if elem.left < 0 or elem.top < 0:
                return True
            if elem.left + elem.width > self.PAGE_WIDTH:
                return True
            if elem.top + elem.height > self.PAGE_HEIGHT:
                return True
        
        return False
    
    def _get_llm_suggestions(self, slide_bboxes: SlideBBoxes) -> List[Dict[str, Any]]:
        """
        获取LLM的布局调整建议（增强版：块级别）
        
        Args:
            slide_bboxes: 单页的bounding boxes（包含block信息）
        
        Returns:
            调整建议列表
        """
        # 构建prompt（美观性导向的prompt）
        prompt = f"""你是一位专业的PPT美学设计师，精通教学型PPT的视觉排版。请基于美观性和用户体验对以下页面布局进行优化。

📐 **页面信息**
- 尺寸：{self.PAGE_WIDTH} x {self.PAGE_HEIGHT} 英寸
- 安全区域：left >= 0.5, top >= 0.5, right <= {self.PAGE_WIDTH - 0.5}, bottom <= {self.PAGE_HEIGHT - 0.5}
- 元素数量：{len(slide_bboxes.elements)}个

📦 **当前元素坐标**（每个元素有4个点：left, top, left+width, top+height）
{json.dumps([e.to_dict() for e in slide_bboxes.elements], ensure_ascii=False, indent=2)}

🎨 **美观性设计原则**（重要！）

1. **黄金分割与视觉平衡**
   - 主要内容应居中或遵循1:1.618比例
   - 避免所有元素都挤在一侧
   - 重要元素（formula, manim_relation）应占据显著位置

2. **留白与呼吸感**
   - 元素之间至少保持0.5英寸的间距
   - 边距至少0.5英寸（不要顶到边界）
   - 文字密集区域需要更多留白

3. **视觉层次**
   - 标题/概念陈述（text类型）：靠上，清晰
   - 动画/图片（image类型）：居中，突出
   - 说明文字：紧贴对应元素，但不遮挡

4. **Manim动画特殊布局**
   - GIF动画：水平居中，垂直居中或偏上
   - 说明文字：紧贴GIF正下方（间距0.1-0.2英寸），水平居中
   - GIF推荐尺寸：width 6-8英寸，height 4-5英寸

5. **无Overlap原则**
   - 任何两个元素都不能重叠（bbox不相交）
   - 如果检测到重叠，必须调整

6. **边界约束（硬性要求）**
   - 所有元素必须在安全区域内
   - left >= 0.5, top >= 0.5
   - left + width <= {self.PAGE_WIDTH - 0.5}
   - top + height <= {self.PAGE_HEIGHT - 0.5}

🔧 **可调整项**
- **text类型**：只能调整left, top（不能改width, height）
- **image类型**：可以调整left, top, width, height
- **manim_relation/manim_visual**：可以调整所有参数

📋 **你的任务**
1. 检查是否有overlap（两个元素的bbox相交）
2. 检查是否有元素超出边界
3. 从美观性角度评估当前布局
4. 给出具体的调整建议（新的left, top, width, height值）

⚠️ **输出格式**（必须是有效的JSON）
{{
  "aesthetic_score": 7,  // 当前美观度评分（0-10）
  "has_issues": true,    // 是否需要调整
  "issues_description": [
    "元素0和元素1存在overlap，影响视觉",
    "元素2超出底部边界0.3英寸",
    "整体布局偏左，缺乏平衡感"
  ],
  "adjustments": [
    {{
      "element_index": 0,
      "reason": "解决与元素1的overlap，同时优化视觉平衡",
      "new_left": 1.5,
      "new_top": 2.0
    }},
    {{
      "element_index": 2,
      "reason": "缩小尺寸以符合边界约束，提升美观度",
      "new_width": 6.0,
      "new_height": 4.0
    }}
  ]
}}

**重要提示**：
- 即使没有严重问题，也可以从美观性角度提出优化建议
- 优先解决overlap和边界溢出
- 注重整体视觉平衡和用户观看体验
- 只返回JSON，不要其他文字"""

        # 重试机制（最多2次）
        import time
        max_retries = 2
        for attempt in range(max_retries):
            try:
                response = self.llm(prompt)
                
                # 检查是否是CustomLLM错误消息
                if "请求大模型失败" in response or "请求大模型超时" in response:
                    if attempt < max_retries - 1:
                        print(f"      ⚠️ API调用失败，重试 {attempt + 1}/{max_retries}")
                        time.sleep(5)  # API失败延迟更长
                        continue
                    else:
                        print(f"      ⚠️ API调用失败: {response[:100]}，跳过本页")
                        return []
                
                # 检查空响应
                if not response or len(response.strip()) < 10:
                    if attempt < max_retries - 1:
                        print(f"      ⚠️ LLM响应为空，重试 {attempt + 1}/{max_retries}")
                        time.sleep(1)
                        continue
                    else:
                        print(f"      ⚠️ LLM响应为空，跳过本页")
                        return []
                
                # 清理响应
                response = response.replace("<br>", "\n")
                response = re.sub(r'<[^>]+>', '', response)  # 移除HTML标签
                response = re.sub(r'```json\s*', '', response, flags=re.IGNORECASE)
                response = re.sub(r'```\s*', '', response)
                response = response.strip()
                
                if not response:
                    if attempt < max_retries - 1:
                        print(f"      ⚠️ 清理后响应为空，重试 {attempt + 1}/{max_retries}")
                        time.sleep(1)
                        continue
                    else:
                        return []
                
                # 解析JSON
                result = json.loads(response)
                
                # 调试：打印完整响应
                print(f"      🔍 LLM返回的完整结果:")
                print(f"         has_issues: {result.get('has_issues')}")
                print(f"         aesthetic_score: {result.get('aesthetic_score')}")
                print(f"         adjustments数量: {len(result.get('adjustments', []))}")
                
                if result.get("has_issues", False):
                    print(f"      发现问题: {result.get('issues_description', [])}")
                    adjustments = result.get("adjustments", [])
                    if not adjustments:
                        print(f"      ⚠️ LLM发现问题但未提供调整方案，返回空调整")
                    return adjustments
                else:
                    return []
                    
            except json.JSONDecodeError as e:
                if attempt < max_retries - 1:
                    print(f"      ⚠️ LLM响应解析失败: {e}，重试 {attempt + 1}/{max_retries}")
                    time.sleep(1)
                else:
                    print(f"      ⚠️ LLM响应解析失败: {e}，跳过本页")
                    return []
            except Exception as e:
                if attempt < max_retries - 1:
                    print(f"      ⚠️ LLM调用失败: {e}，重试 {attempt + 1}/{max_retries}")
                    time.sleep(1)
                else:
                    print(f"      ⚠️ LLM调用失败: {e}，跳过本页")
                    return []
        
        return []
    
    def _apply_adjustments(self, slide_idx: int, adjustments: List[Dict[str, Any]], slide_bboxes: SlideBBoxes):
        """
        应用布局调整（增强版：使用shape_index）
        
        Args:
            slide_idx: slide索引
            adjustments: 调整列表（包含block_id和reason）
            slide_bboxes: 该页的bounding boxes（包含shape_index）
        """
        slide = self.prs.slides[slide_idx]
        shapes = list(slide.shapes)
        
        for adj in adjustments:
            try:
                elem_idx = adj.get("element_index")  # BBox列表中的索引
                block_id = adj.get("block_id", "unknown")
                reason = adj.get("reason", "")
                
                if elem_idx is None or elem_idx >= len(slide_bboxes.elements):
                    continue
                
                # 关键：使用shape_index而不是element_index
                bbox = slide_bboxes.elements[elem_idx]
                shape_idx = bbox.shape_index if bbox.shape_index is not None else elem_idx
                
                if shape_idx >= len(shapes):
                    print(f"      ⚠️ shape_index {shape_idx} 超出范围 (总共{len(shapes)}个元素)")
                    continue
                
                shape = shapes[shape_idx]
                
                # 自动推断action类型（兼容无action字段的LLM输出）
                has_position = "new_left" in adj or "new_top" in adj
                has_size = "new_width" in adj or "new_height" in adj
                
                # 应用位置调整
                if has_position:
                    new_left = adj.get("new_left")
                    new_top = adj.get("new_top")
                    
                    if new_left is not None:
                        shape.left = Inches(new_left)
                    if new_top is not None:
                        shape.top = Inches(new_top)
                    
                    self.adjustment_log.append({
                        "slide_idx": slide_idx,
                        "element_idx": elem_idx,
                        "shape_idx": shape_idx,  # 记录实际的shape索引
                        "block_id": block_id,
                        "action": "move",
                        "reason": reason,
                        "new_left": new_left,
                        "new_top": new_top
                    })
                    
                    print(f"      ✓ [Block {block_id}] 移动元素 {elem_idx} (shape {shape_idx}): ({new_left}, {new_top})")
                    if reason:
                        print(f"        原因: {reason}")
                
                # 应用尺寸调整
                if has_size:
                    new_width = adj.get("new_width")
                    new_height = adj.get("new_height")
                    
                    if new_width is not None:
                        shape.width = Inches(new_width)
                    if new_height is not None:
                        shape.height = Inches(new_height)
                    
                    self.adjustment_log.append({
                        "slide_idx": slide_idx,
                        "element_idx": elem_idx,
                        "shape_idx": shape_idx,  # 记录实际的shape索引
                        "block_id": block_id,
                        "action": "resize",
                        "reason": reason,
                        "new_width": new_width,
                        "new_height": new_height
                    })
                    
                    print(f"      ✓ [Block {block_id}] 调整大小 {elem_idx} (shape {shape_idx}): {new_width}x{new_height}")
                    if reason:
                        print(f"        原因: {reason}")
                
            except Exception as e:
                print(f"      ⚠️ 应用调整失败 (element {elem_idx}, shape {shape_idx if 'shape_idx' in locals() else 'unknown'}): {e}")
    
    def _save_adjustment_log(self):
        """保存调整日志"""
        log_path = self.output_dir / "layout_adjustments.json"
        
        with open(log_path, 'w', encoding='utf-8') as f:
            json.dump({
                "pptx_file": str(self.pptx_path),
                "adjustments": self.adjustment_log
            }, f, ensure_ascii=False, indent=2)
        
        print(f"  📝 调整日志已保存: {log_path}")


def optimize_layout(pptx_path: str, bboxes: List[SlideBBoxes], output_dir: str) -> Path:
    """
    便捷函数：优化PPTX布局
    
    Args:
        pptx_path: PPTX文件路径
        bboxes: bounding boxes列表
        output_dir: 输出目录
    
    Returns:
        优化后的PPTX路径
    """
    optimizer = VLMLayoutOptimizer(pptx_path, output_dir)
    return optimizer.optimize(bboxes)


