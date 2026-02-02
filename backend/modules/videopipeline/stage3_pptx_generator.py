# stage3_pptx_generator.py
# Stage 3: PPTX生成器

import os
import re
from typing import List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

from .models import SlideStructure, SlideType


class PPTXGenerator:
    """
    PPTX生成器
    
    功能：
    1. 使用模板创建PPTX（不修改模板结构）
    2. 第一页是封面页（只修改标题，删除二级标题和日期）
    3. 后续内容页使用第二页（空白页）作为模板
    
    字体格式参考 lecture_generator.py
    """
    
    # PPT尺寸（16:9）
    SLIDE_WIDTH = 13.33  # 英寸
    SLIDE_HEIGHT = 7.5   # 英寸
    
    # 布局常量
    TITLE_TOP = 0.3
    TITLE_HEIGHT = 0.8
    CONTENT_TOP = 1.3
    CONTENT_LEFT = 0.5
    TEXT_WIDTH = 6.0
    TEXT_HEIGHT = 5.5
    GIF_LEFT = 7.0
    GIF_TOP = 1.5
    GIF_WIDTH = 5.5
    GIF_HEIGHT = 5.0
    
    def __init__(self, template_path: str, output_dir: str):
        """
        初始化生成器
        
        Args:
            template_path: PPTX模板路径
            output_dir: 输出目录
        """
        self.template_path = Path(template_path)
        self.output_dir = Path(output_dir)
        self.pptx_dir = self.output_dir / "pptx"
        self.pptx_dir.mkdir(parents=True, exist_ok=True)
        
        self.prs: Optional[Presentation] = None
        self.blank_layout = None  # 第二页的空白布局
    
    def generate(self, slides: List[SlideStructure], title: str = "教学课件") -> Path:
        """
        生成PPTX文件
        
        Args:
            slides: slide结构列表
            title: 课件标题（用于修改第一页）
        
        Returns:
            生成的PPTX文件路径
        """
        # 加载模板
        if not self.template_path.exists():
            raise FileNotFoundError(f"模板文件不存在: {self.template_path}")
        
        self.prs = Presentation(str(self.template_path))
        print(f"  📄 加载模板: {self.template_path}")
        print(f"     模板共 {len(self.prs.slides)} 页")
        
        # 修改第一页（封面页）：只改标题，删除二级标题和日期
        self._modify_cover_slide(title)
        
        # 删除第二页（模板空白页），保存其布局以供后续使用
        # 注意：第二页是索引1
        if len(self.prs.slides) >= 2:
            # 记录第二页的布局（用于添加新页）
            self.blank_layout = self.prs.slides[1].slide_layout
            # 删除第二页
            rId = self.prs.slides._sldIdLst[1].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[1]
        else:
            # 如果模板只有一页，使用空白布局
            self.blank_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[0]
        
        # 添加内容页
        for slide in slides:
            self._add_content_slide(slide)
        
        # 保存文件
        output_path = self.pptx_dir / "presentation.pptx"
        self.prs.save(str(output_path))
        print(f"  ✅ PPTX已保存: {output_path}")
        
        return output_path
    
    def _modify_cover_slide(self, title: str):
        """
        修改封面页：只修改主标题，删除二级标题和日期
        
        Args:
            title: 新的主标题
        """
        if len(self.prs.slides) == 0:
            return
        
        cover_slide = self.prs.slides[0]
        shapes_to_delete = []
        title_shape = None
        
        # 找到所有文本框并按面积排序，最大的是主标题
        text_shapes = []
        for shape in cover_slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                area = shape.width * shape.height
                text_shapes.append((area, shape))
        
        # 按面积排序
        text_shapes.sort(key=lambda x: x[0], reverse=True)
        
        if text_shapes:
            # 最大的文本框是主标题
            title_shape = text_shapes[0][1]
            
            # 其他所有文本框都删除（二级标题、日期等）
            for _, shape in text_shapes[1:]:
                shapes_to_delete.append(shape)
        
        # 修改主标题
        if title_shape:
            # 清除所有段落文本
            for para in title_shape.text_frame.paragraphs:
                para.clear()
            # 设置新标题
            title_shape.text_frame.paragraphs[0].text = title
            for run in title_shape.text_frame.paragraphs[0].runs:
                run.font.name = '微软雅黑'
                run.font.bold = True
                run.font.size = Pt(44)
                run.font.color.rgb = RGBColor(192, 0, 0)
        
        # 删除二级标题和日期
        for shape in shapes_to_delete:
            try:
                sp = shape.element
                sp.getparent().remove(sp)
            except Exception as e:
                print(f"      ⚠️ 无法删除元素: {e}")
    
    def _add_content_slide(self, slide_data: SlideStructure):
        """
        添加内容页
        
        使用第二页（空白页）的布局，不修改模板结构
        标题使用section_title（从test.json的script.sections[].title获取）
        
        每个内容块（Block）都生成为一个独立的 Shape，以便后续 VLM 优化布局
        并添加“出现”动画，以便视频合成时实现块级对齐
        
        Args:
            slide_data: slide结构体
        """
        # 使用保存的空白布局（来自模板第二页）
        slide = self.prs.slides.add_slide(self.blank_layout)
        
        # 标题优先使用section_title（来自test.json）
        display_title = slide_data.section_title if slide_data.section_title else slide_data.title
        self._add_title(slide, display_title)
        
        # 记录当前布局的纵向位置
        current_top = self.CONTENT_TOP
        
        # 如果没有 blocks（回退情况），使用旧逻辑
        if not slide_data.blocks:
            if slide_data.slide_type == SlideType.FORMULA and slide_data.formula:
                self._add_text_with_formula(slide, slide_data)
            elif slide_data.gif_path and slide_data.gif_path.exists():
                self._add_text_with_gif(slide, slide_data)
            else:
                self._add_text_full_width(slide, slide_data)
            return

        # 遍历所有 blocks 并生成独立的 Shape
        for i, block in enumerate(slide_data.blocks):
            shape = None
            block_type_str = block.block_type.value if hasattr(block.block_type, 'value') else str(block.block_type)
            
            if block_type_str in ["text_line", "conceptual_statement"]:
                # 文本类型
                shape = self._add_block_text(slide, block, current_top, slide_data)
                if shape:
                    current_top += shape.height.inches + 0.2
            elif block_type_str == "formula":
                shape = self._add_block_formula(slide, block, current_top)
                if shape:
                    current_top += shape.height.inches + 0.2
            elif block_type_str == "code":
                # 代码块 - 优先使用GIF（如果存在）
                gif_path = slide_data.gif_path if hasattr(slide_data, 'gif_path') and slide_data.gif_path else None
                
                if gif_path and Path(gif_path).exists():
                    # 使用Manim GIF渲染代码
                    shape = self._add_block_gif(slide, Path(gif_path))
                    if shape:
                        current_top += shape.height.inches + 0.3
                else:
                    # 回退到文本渲染
                    shape = self._add_block_code(slide, block, current_top)
                    if shape:
                        current_top += shape.height.inches + 0.2
            elif block_type_str in ["manim_relation", "manim_visual"]:
                # Manim 关系图 - 先添加GIF，然后在GIF下方添加说明文字（带点击效果）
                gif_path = None
                ppt_text = None
                
                if isinstance(block.content, dict):
                    gif_path = Path(block.content['gif_path']) if 'gif_path' in block.content else None
                    ppt_text = block.content.get('ppt_text', block.content.get('description', ''))
                elif hasattr(block, 'manim_code') and block.manim_code:
                    gif_path = slide_data.gif_path
                    ppt_text = "关系图"
                
                # 添加GIF（居中）
                gif_shape = None
                text_shape = None
                if gif_path and gif_path.exists():
                    gif_shape = self._add_block_gif_centered(slide, gif_path, current_top)
                    if gif_shape:
                        current_top = gif_shape.top.inches + gif_shape.height.inches + 0.1
                
                # 添加说明文字（紧贴GIF下方，居中）
                if ppt_text and gif_shape:
                    text_shape = self._add_manim_text_below(slide, ppt_text, current_top)
                    if text_shape:
                        current_top += text_shape.height.inches + 0.3
                        shape = text_shape  # 记录最后一个shape
                elif gif_shape:
                    shape = gif_shape
                    current_top += 0.3
                elif not gif_path:
                    # 如果没有GIF，显示占位符
                    shape = self._add_block_placeholder(slide, "[关系图]", current_top)
                    current_top += shape.height.inches + 0.2
                
            elif block_type_str == "image":
                # 图片 - 从 URL 下载或从本地路径插入（检查边界）
                shape = self._add_block_image_bounded(slide, block, current_top)
                if shape:
                    current_top += shape.height.inches + 0.2
            
            # 为每个 block 添加"出现"动画（Manim blocks已单独处理）
            if shape and block_type_str not in ["manim_relation", "manim_visual"]:
                self._add_entrance_animation(shape)

    def _add_block_text(self, slide, block, top, slide_data=None):
        """
        添加文本块（支持emphasis样式）
        
        实现逐词处理，支持加粗和颜色
        """
        # 检查slide_data.blocks是否有manim/image类型（用于决定文本宽度）
        has_visual = False
        if slide_data and hasattr(slide_data, 'blocks'):
            for b in slide_data.blocks:
                block_type = b.block_type.value if hasattr(b.block_type, 'value') else str(b.block_type)
                if block_type in ["manim_visual", "image", "manim_relation"]:
                    has_visual = True
                    break
        
        width = self.TEXT_WIDTH if has_visual else self.SLIDE_WIDTH - 1.0
        
        text_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(top),
            Inches(width),
            Inches(0.5)  # 初始高度，会自动扩展
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        # 如果没有emphasis，直接设置文本
        if not block.emphasis or (not block.emphasis.get('bold') and not block.emphasis.get('color')):
            p = tf.paragraphs[0]
            p.text = block.content
            p.font.name = '微软雅黑'
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0, 0, 0)
            return text_box
        
        # 有emphasis：逐词处理
        self._apply_text_emphasis(tf, block.content, block.emphasis)
        
        return text_box
    
    def _apply_text_emphasis(self, text_frame, content: str, emphasis: dict):
        """
        应用文本强调样式（加粗+颜色）
        
        策略：按关键词分割文本，为每个部分创建独立的run
        """
        bold_words = set(emphasis.get('bold', []))
        color_map = emphasis.get('color', {})
        
        para = text_frame.paragraphs[0]
        
        # 构建关键词列表（按长度排序，优先匹配长词）
        keywords = list(bold_words | set(color_map.keys()))
        keywords.sort(key=len, reverse=True)
        
        if not keywords:
            # 没有关键词，直接设置
            para.text = content
            self._format_text_run(para.runs[0], bold=False, color=None)
            return
        
        # 使用正则表达式分割文本
        import re
        pattern = '(' + '|'.join(re.escape(kw) for kw in keywords) + ')'
        parts = re.split(pattern, content)
        
        first = True
        for part in parts:
            if not part:
                continue
            
            # 判断是否是关键词
            is_bold = part in bold_words
            keyword_color = color_map.get(part)
            
            # 添加run
            if first:
                run = para.runs[0] if para.runs else para.add_run()
                first = False
            else:
                run = para.add_run()
            
            run.text = part
            self._format_text_run(run, bold=is_bold, color=keyword_color)
    
    def _format_text_run(self, run, bold=False, color=None):
        """格式化单个text run"""
        run.font.name = '微软雅黑'
        run.font.size = Pt(18)
        
        if bold:
            run.font.bold = True
        
        if color:
            color_rgb = self._parse_color(color)
            run.font.color.rgb = color_rgb
        else:
            run.font.color.rgb = RGBColor(0, 0, 0)
    
    def _parse_color(self, color_name: str) -> RGBColor:
        """解析颜色名称为RGB"""
        color_map = {
            "red": RGBColor(192, 0, 0),      # 深红色
            "blue": RGBColor(0, 0, 128),     # 深蓝色
            "green": RGBColor(0, 100, 0),    # 深绿色
            "orange": RGBColor(255, 140, 0), # 橙色
            "purple": RGBColor(128, 0, 128), # 紫色
        }
        return color_map.get(color_name, RGBColor(0, 0, 0))

    def _add_block_formula(self, slide, block, top):
        """添加公式块"""
        formula_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT + 0.5),
            Inches(top),
            Inches(self.SLIDE_WIDTH - 2.0),
            Inches(0.8)
        )
        tf = formula_box.text_frame
        p = tf.paragraphs[0]
        p.text = block.content
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Cambria Math'
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 0, 128)
        return formula_box

    def _add_block_code(self, slide, block, top):
        """添加代码块"""
        code_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(top),
            Inches(self.SLIDE_WIDTH - 1.0),
            Inches(2.0)
        )
        tf = code_box.text_frame
        tf.text = block.content
        for para in tf.paragraphs:
            para.font.name = 'Consolas'
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(0, 100, 0)
        return code_box

    def _add_block_gif(self, slide, gif_path):
        """添加 GIF 块（右侧位置）"""
        return slide.shapes.add_picture(
            str(gif_path),
            Inches(7.5),  # 右侧
            Inches(2.0),  # 顶部留空
            width=Inches(5.0),  # 固定宽度
        )
    
    def _add_block_gif_centered(self, slide, gif_path, top):
        """添加 GIF 块（居中位置，防止超出边界）"""
        # 设置最大尺寸，防止超出PPT边界
        max_width = Inches(8.0)  # 留出左右边距
        max_height = Inches(4.0)  # 留出上下边距
        
        # 添加图片（python-pptx会自动保持比例）
        pic = slide.shapes.add_picture(
            str(gif_path),
            Inches(1.5),  # 左侧留白
            Inches(top),
            width=max_width
        )
        
        # 如果高度超出，重新调整
        if pic.height > max_height:
            pic.height = max_height
            # 重新居中
            pic.left = Inches((13.33 - pic.width.inches) / 2)
        else:
            # 水平居中
            pic.left = Inches((13.33 - pic.width.inches) / 2)
        
        return pic
    
    def _add_manim_text(self, slide, text, top):
        """添加Manim说明文字（左侧位置）"""
        text_box = slide.shapes.add_textbox(
            Inches(0.5),  # 左侧
            Inches(top),
            Inches(6.5),  # 宽度（不与右侧GIF重叠）
            Inches(1.5)
        )
        
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = text
        for para in tf.paragraphs:
            para.font.size = Pt(16)
            para.font.bold = False
            para.font.color.rgb = RGBColor(50, 50, 50)
            para.alignment = 1  # 左对齐
        
        return text_box
    
    def _add_manim_text_below(self, slide, text, top):
        """添加Manim说明文字（紧贴GIF下方，居中）"""
        text_box = slide.shapes.add_textbox(
            Inches(2.0),  # 左边距
            Inches(top),
            Inches(9.33),  # 居中宽度
            Inches(1.0)
        )
        
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = text
        
        for para in tf.paragraphs:
            para.font.size = Pt(14)
            para.font.bold = False
            para.font.color.rgb = RGBColor(80, 80, 80)
            para.alignment = 1  # 居中对齐（PP_ALIGN.CENTER = 1）
        
        return text_box
    
    def _add_block_image_bounded(self, slide, block, top):
        """添加图片块（支持URL或本地路径，防止超出边界）"""
        import requests
        import tempfile
        from pathlib import Path
        
        image_path = None
        is_temp = False
        
        # 获取图片路径或URL
        image_source = block.content if isinstance(block.content, str) else block.image_url
        
        if not image_source:
            return self._add_block_placeholder(slide, "[Image placeholder]", top)
        
        # 判断是URL还是本地路径
        if image_source.startswith('http://') or image_source.startswith('https://'):
            # 从URL下载图片
            try:
                response = requests.get(image_source, timeout=10)
                response.raise_for_status()
                
                # 保存到临时文件
                suffix = '.jpg'  # 默认jpg
                if 'image/png' in response.headers.get('Content-Type', ''):
                    suffix = '.png'
                elif 'image/gif' in response.headers.get('Content-Type', ''):
                    suffix = '.gif'
                
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
                temp_file.write(response.content)
                temp_file.close()
                
                image_path = temp_file.name
                is_temp = True
                
            except Exception as e:
                print(f"      ⚠️ 下载图片失败: {e}")
                return self._add_block_placeholder(slide, f"[Image: {image_source[:50]}...]", top)
        else:
            # 本地路径
            image_path = Path(image_source)
            if not image_path.exists():
                return self._add_block_placeholder(slide, f"[Image not found: {image_source}]", top)
        
        # 插入图片到PPT（右侧位置，防止超出边界）
        try:
            # 设置最大尺寸，确保不超出PPT边界
            max_width = Inches(4.5)
            max_height = Inches(5.0)  # 避免超出底部边界
            
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(8.0),  # 右侧
                Inches(top),
                width=max_width
            )
            
            # 如果高度超出，调整大小
            if picture.height > max_height:
                picture.height = max_height
            
            # 确保不超出底部边界（PPT高度约为7.5英寸）
            if picture.top.inches + picture.height.inches > 7.0:
                new_height = Inches(7.0 - picture.top.inches)
                picture.height = new_height
            
            # 清理临时文件
            if is_temp:
                try:
                    Path(image_path).unlink()
                except:
                    pass
            
            return picture
            
        except Exception as e:
            print(f"      ⚠️ 插入图片失败: {e}")
            if is_temp:
                try:
                    Path(image_path).unlink()
                except:
                    pass
            return self._add_block_placeholder(slide, "[Image insert failed]", top)

    def _add_block_placeholder(self, slide, text, top):
        """添加占位块"""
        box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(top),
            Inches(3.0),
            Inches(1.0)
        )
        box.text_frame.text = text
        return box

    def _add_entrance_animation(self, shape, effect="Appear"):
        """
        为形状添加简单的进入动画（点击触发）
        
        Args:
            shape: 要添加动画的形状
            effect: 动画效果类型（默认为"Appear"）
        
        Note:
            python-pptx对动画支持有限，这里主要是设置隐藏初始状态
            真正的动画效果由Stage 8在视频合成时控制
        """
        try:
            # python-pptx没有直接的动画API，但我们可以通过设置形状的可见性
            # 在Stage 8中，我们会根据timing信息控制每个block的出现时机
            # 这里只是标记，实际动画在视频合成阶段处理
            pass
        except Exception as e:
            print(f"      ⚠️ 添加动画标记失败: {e}")

    def _add_title(self, slide, title: str):
        """添加标题"""
        title_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(self.TITLE_TOP),
            Inches(self.SLIDE_WIDTH - 1),
            Inches(self.TITLE_HEIGHT)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        
        for para in title_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.name = '微软雅黑'
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(192, 0, 0)  # 深红色
    
    def _add_text_with_gif(self, slide, slide_data: SlideStructure):
        """添加带GIF的布局：左文字右GIF"""
        # 左侧文本框
        text_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(self.CONTENT_TOP),
            Inches(self.TEXT_WIDTH),
            Inches(self.TEXT_HEIGHT)
        )
        self._format_text_box(text_box, slide_data.text)
        
        # 右侧GIF
        slide.shapes.add_picture(
            str(slide_data.gif_path),
            Inches(self.GIF_LEFT),
            Inches(self.GIF_TOP),
            Inches(self.GIF_WIDTH),
            Inches(self.GIF_HEIGHT)
        )
    
    def _add_text_with_formula(self, slide, slide_data: SlideStructure):
        """添加带公式的布局：左文字右公式文本"""
        # 左侧文本框
        text_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(self.CONTENT_TOP),
            Inches(self.TEXT_WIDTH),
            Inches(self.TEXT_HEIGHT)
        )
        self._format_text_box(text_box, slide_data.text)
        
        # 右侧公式文本框（使用Cambria Math字体）
        self._add_formula_text_box(slide, slide_data.formula)
    
    def _add_text_with_code(self, slide, slide_data: SlideStructure):
        """添加带代码的布局：左文字右代码文本"""
        # 左侧文本框
        text_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(self.CONTENT_TOP),
            Inches(self.TEXT_WIDTH),
            Inches(self.TEXT_HEIGHT)
        )
        self._format_text_box(text_box, slide_data.text)
        
        # 右侧代码文本框
        self._add_code_text_box(slide, slide_data.coq_code)
    
    def _add_text_full_width(self, slide, slide_data: SlideStructure):
        """添加全幅文字布局"""
        text_box = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT),
            Inches(self.CONTENT_TOP),
            Inches(self.SLIDE_WIDTH - 1),
            Inches(self.TEXT_HEIGHT)
        )
        self._format_text_box(text_box, slide_data.text)
        
        # 如果有代码但没有GIF，显示代码文本
        if slide_data.coq_code:
            self._add_code_text_box(slide, slide_data.coq_code)
        
        # 如果有公式但没有GIF，显示公式文本
        if slide_data.formula:
            self._add_formula_text_box(slide, slide_data.formula)
    
    def _format_text_box(self, text_box, text: str):
        """
        格式化文本框，使用普通bullet point格式
        
        不使用模板的logo，使用简单的文字要点格式
        """
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        # 分段处理 - 每行作为一个要点
        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        
        for i, para_text in enumerate(paragraphs):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()
            
            # 去除可能的bullet符号
            para_text = re.sub(r'^[•\-\*\d\.]+\s*', '', para_text)
            
            # 如果有多个要点，添加简单的bullet符号
            if len(paragraphs) > 1:
                para.text = "• " + para_text
            else:
                para.text = para_text
            
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(12)
            para.level = 0  # 不使用缩进级别（避免触发模板的bullet样式）
            
            for run in para.runs:
                run.font.name = '微软雅黑'
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0, 0, 0)
    
    def _add_code_text_box(self, slide, code: str):
        """添加代码文本框（当没有GIF时使用）"""
        code_box = slide.shapes.add_textbox(
            Inches(self.GIF_LEFT),
            Inches(self.GIF_TOP),
            Inches(self.GIF_WIDTH),
            Inches(self.GIF_HEIGHT)
        )
        
        text_frame = code_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = code
        
        # 设置代码字体
        for para in text_frame.paragraphs:
            for run in para.runs:
                run.font.name = 'Consolas'
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0, 100, 0)  # 深绿色
        
        # 添加背景色（通过形状）
        # 注意：python-pptx不直接支持文本框背景色
        # 可以在文本框下方添加一个填充矩形作为背景
    
    def _add_formula_text_box(self, slide, formula: str):
        """添加公式文本框（当没有GIF时使用）"""
        formula_box = slide.shapes.add_textbox(
            Inches(self.GIF_LEFT),
            Inches(self.GIF_TOP),
            Inches(self.GIF_WIDTH),
            Inches(2.0)
        )
        
        text_frame = formula_box.text_frame
        text_frame.text = formula
        
        for para in text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.name = 'Cambria Math'
                run.font.size = Pt(24)
                run.font.color.rgb = RGBColor(0, 0, 128)  # 深蓝色


def generate_pptx(
    slides: List[SlideStructure],
    template_path: str,
    output_dir: str,
    title: str = "教学课件"
) -> Path:
    """
    便捷函数：生成PPTX
    
    Args:
        slides: slide列表
        template_path: 模板路径
        output_dir: 输出目录
        title: 课件标题
    
    Returns:
        生成的PPTX文件路径
    """
    generator = PPTXGenerator(template_path, output_dir)
    return generator.generate(slides, title)


