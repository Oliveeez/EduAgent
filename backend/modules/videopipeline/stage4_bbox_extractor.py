# stage4_bbox_extractor.py
# Stage 4: Bounding Box提取（块级别）

import json
from typing import List, Dict, Any, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from .models import BoundingBox, SlideBBoxes, SlideStructure


class BBoxExtractor:
    """
    Bounding Box提取器（增强版：块级别）
    
    功能：
    1. 从PPTX中提取所有元素的位置和大小
    2. 将元素与blocks关联（通过block_id）
    3. 生成块级别的布局分析JSON
    4. 检测块之间的重叠问题
    """
    
    def __init__(self, pptx_path: str, slides_data: Optional[List[SlideStructure]] = None):
        """
        初始化提取器
        
        Args:
            pptx_path: PPTX文件路径
            slides_data: slide结构数据（包含blocks信息）
        """
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(self.pptx_path))
        self.slides_data = slides_data or []
        self.all_bboxes: List[SlideBBoxes] = []
    
    def extract(self) -> List[SlideBBoxes]:
        """
        提取所有slides的bounding boxes
        
        Returns:
            每页的bounding boxes列表
        """
        self.all_bboxes = []
        
        for slide_idx, slide in enumerate(self.prs.slides):
            slide_bboxes = self._extract_slide_bboxes(slide, slide_idx)
            self.all_bboxes.append(slide_bboxes)
        
        return self.all_bboxes
    
    def _extract_slide_bboxes(self, slide, slide_idx: int) -> SlideBBoxes:
        """
        提取单页的bounding boxes（块级别，过滤空元素）
        
        Args:
            slide: pptx slide对象
            slide_idx: slide索引
        
        Returns:
            该页的SlideBBoxes（只包含有内容的元素，关联到blocks）
        """
        slide_bboxes = SlideBBoxes(slide_id=slide_idx, elements=[])
        
        # 获取对应的slide_data（如果有）
        slide_data = None
        if slide_idx < len(self.slides_data):
            slide_data = self.slides_data[slide_idx]
        
        # 提取所有shapes（无论是否匹配到block）
        shapes_list = list(slide.shapes)
        
        # 跳过标题slide（slide_placeholder类型）
        for shape_idx, shape in enumerate(shapes_list):
            # 跳过占位符
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                continue
            
            # 提取bbox（记录实际shape索引）
            bbox = self._extract_shape_bbox(shape, None, None)
            if bbox:
                bbox.shape_index = shape_idx  # 关键：记录PPT中的实际索引
                slide_bboxes.elements.append(bbox)
        
        # 如果有slide_data，尝试将提取的bbox与block关联
        if slide_data and hasattr(slide_data, 'blocks') and slide_data.blocks:
            self._associate_bboxes_with_blocks(slide_bboxes.elements, slide_data.blocks)
        
        return slide_bboxes
    
    def _associate_bboxes_with_blocks(self, bboxes: List[BoundingBox], blocks: list):
        """
        将已提取的bboxes与blocks关联
        
        Args:
            bboxes: 已提取的bbox列表
            blocks: SlideBlock列表
        """
        # 简单策略：按顺序关联
        for bbox_idx, bbox in enumerate(bboxes):
            if bbox_idx < len(blocks):
                block = blocks[bbox_idx]
                bbox.block_id = block.block_id
                block_type = block.block_type.value if hasattr(block.block_type, 'value') else str(block.block_type)
                bbox.block_type = block_type
    
    def _extract_slide_bboxes_old(self, slide, slide_idx: int) -> SlideBBoxes:
        """
        旧版提取逻辑（保留作为备用）
        """
        slide_bboxes = SlideBBoxes(slide_id=slide_idx, elements=[])
        
        # 获取对应的slide_data（如果有）
        slide_data = None
        if slide_idx < len(self.slides_data):
            slide_data = self.slides_data[slide_idx]
        
        # 提取所有shapes
        shapes_list = list(slide.shapes)
        
        # 如果有slide_data，尝试将shape与block匹配
        if slide_data and hasattr(slide_data, 'blocks') and slide_data.blocks:
            block_shape_map = self._match_shapes_to_blocks(shapes_list, slide_data.blocks)
            
            # 按照block顺序提取bbox
            for block in slide_data.blocks:
                shape_indices = block_shape_map.get(block.block_id, [])
                for shape_idx in shape_indices:
                    if shape_idx < len(shapes_list):
                        shape = shapes_list[shape_idx]
                        bbox = self._extract_shape_bbox(
                            shape, 
                            block_id=block.block_id,
                            block_type=block.block_type.value if hasattr(block.block_type, 'value') else str(block.block_type),
                            semantic_role=block.semantic_role.value if block.semantic_role and hasattr(block.semantic_role, 'value') else None
                        )
                        if bbox:
                            slide_bboxes.elements.append(bbox)
        else:
            # 没有block信息，按原来的方式提取
            for shape in shapes_list:
                bbox = self._extract_shape_bbox(shape)
                if bbox:
                    slide_bboxes.elements.append(bbox)
        
        return slide_bboxes
    
    def _match_shapes_to_blocks(self, shapes: list, blocks: list) -> Dict[str, List[int]]:
        """
        将shapes与blocks匹配
        
        Args:
            shapes: pptx shapes列表
            blocks: SlideBlock列表
        
        Returns:
            {block_id: [shape_index, ...]}
        """
        block_shape_map = {}
        
        # 简单的匹配策略：按照顺序匹配
        # 假设PPT生成时，block和shape的顺序是一致的
        shape_idx = 0
        for block in blocks:
            block_shape_map[block.block_id] = []
            
            # 每个block可能对应1个或多个shape
            # 对于text block，通常是1个textbox
            # 对于manim_visual block，可能是1个图片
            # 对于code block，可能是多个textbox
            
            if shape_idx < len(shapes):
                shape = shapes[shape_idx]
                
                # 检查shape内容是否与block匹配
                if hasattr(shape, 'text_frame') and shape.text_frame.text:
                    text_content = shape.text_frame.text.strip()
                    block_content = str(block.content).strip()
                    
                    # 如果文本内容匹配或包含，则认为是同一个block
                    if text_content in block_content or block_content in text_content:
                        block_shape_map[block.block_id].append(shape_idx)
                        shape_idx += 1
                    elif len(text_content) > 0:
                        # 有文本但不匹配，也加入（可能是格式化后的文本）
                        block_shape_map[block.block_id].append(shape_idx)
                        shape_idx += 1
                elif hasattr(shape, 'image'):
                    # 图片/GIF
                    block_shape_map[block.block_id].append(shape_idx)
                    shape_idx += 1
        
        return block_shape_map
    
    def _calculate_text_height(self, shape) -> float:
        """
        计算文本实际渲染高度（考虑行数、字体大小、行距）
        
        Args:
            shape: 带有text_frame的shape对象
        
        Returns:
            实际文本高度（英寸）
        """
        try:
            text_frame = shape.text_frame
            
            # 1. 统计总行数
            total_lines = 0
            max_font_size = 18  # 默认字体大小（pt）
            
            for paragraph in text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                
                # 获取字体大小
                if paragraph.runs:
                    for run in paragraph.runs:
                        if run.font.size:
                            font_size_pt = run.font.size.pt
                            max_font_size = max(max_font_size, font_size_pt)
                
                # 计算该段落的行数（考虑自动换行）
                # 简化：假设每行最多容纳一定字符数
                text_width = shape.width.inches if hasattr(shape.width, 'inches') else shape.width / 914400
                chars_per_line = int(text_width / 0.15)  # 粗略估计：0.15英寸/字符
                lines_in_para = max(1, len(text) // chars_per_line + (1 if len(text) % chars_per_line else 0))
                total_lines += lines_in_para
            
            if total_lines == 0:
                return shape.height.inches if hasattr(shape.height, 'inches') else shape.height / 914400
            
            # 2. 计算实际高度
            # 行高 = 字体大小 * 行距系数（通常为1.2-1.5）
            line_height_inches = (max_font_size / 72.0) * 1.3  # 72 pt/inch, 1.3倍行距
            
            # 3. 添加padding（上下各0.1英寸）
            padding = 0.2
            
            actual_height = total_lines * line_height_inches + padding
            
            return actual_height
            
        except Exception as e:
            # 失败时返回shape原始高度
            return shape.height.inches if hasattr(shape.height, 'inches') else shape.height / 914400
    
    def _extract_shape_bbox(
        self, 
        shape, 
        block_id: Optional[str] = None, 
        block_type: Optional[str] = None,
        semantic_role: Optional[str] = None
    ) -> Optional[BoundingBox]:
        """
        提取单个形状的bounding box（增强版：关联block）
        
        Args:
            shape: pptx shape对象
            block_id: 关联的block ID
            block_type: block类型
            semantic_role: 语义角色
        
        Returns:
            BoundingBox对象，如果是空元素则返回None
        """
        # 确定元素类型和内容
        element_type = None
        content = ""
        
        if hasattr(shape, 'text_frame'):
            element_type = "text"
            content = shape.text_frame.text.strip() if shape.text_frame.text else ""
            # 过滤空文本框（无内容或只有空白）
            if not content:
                return None
            
            # 计算文本实际渲染高度（考虑字体大小、行数、行距）
            actual_text_height = self._calculate_text_height(shape)
        elif hasattr(shape, 'image'):
            element_type = "image"
            content = "GIF/图片"
        else:
            # 其他形状类型，如果没有内容则跳过
            return None
        
        # 提取位置和大小
        try:
            # 对于文本，使用实际渲染高度；对于图片，使用shape高度
            if element_type == "text" and 'actual_text_height' in locals():
                height = max(
                    actual_text_height,
                    shape.height.inches if hasattr(shape.height, 'inches') else shape.height / 914400
                )
            else:
                height = shape.height.inches if hasattr(shape.height, 'inches') else shape.height / 914400
            
            bbox = BoundingBox(
                element_type=element_type,
                left=shape.left.inches if hasattr(shape.left, 'inches') else shape.left / 914400,
                top=shape.top.inches if hasattr(shape.top, 'inches') else shape.top / 914400,
                width=shape.width.inches if hasattr(shape.width, 'inches') else shape.width / 914400,
                height=height,  # 使用计算后的高度
                content=content[:100],  # 限制长度
                block_id=block_id,
                block_type=block_type,
                semantic_role=semantic_role
            )
            return bbox
        except Exception as e:
            print(f"  ⚠️ 提取bbox失败: {e}")
            return None
    
    def find_overlaps(self) -> List[Dict[str, Any]]:
        """
        查找所有重叠问题
        
        Returns:
            重叠问题列表
        """
        overlaps = []
        
        for slide_bboxes in self.all_bboxes:
            if slide_bboxes.has_overlaps():
                # 找出具体哪些元素重叠
                elements = slide_bboxes.elements
                for i, elem1 in enumerate(elements):
                    for j, elem2 in enumerate(elements[i+1:], i+1):
                        if elem1.overlaps_with(elem2):
                            overlaps.append({
                                "slide_id": slide_bboxes.slide_id,
                                "element1_index": i,
                                "element1_type": elem1.element_type,
                                "element2_index": j,
                                "element2_type": elem2.element_type
                            })
        
        return overlaps
    
    def check_boundaries(self, page_width: float = 13.33, page_height: float = 7.5) -> List[Dict[str, Any]]:
        """
        检查是否有元素超出页面边界
        
        Args:
            page_width: 页面宽度（英寸）
            page_height: 页面高度（英寸）
        
        Returns:
            超出边界的元素列表
        """
        out_of_bounds = []
        
        for slide_bboxes in self.all_bboxes:
            for i, elem in enumerate(slide_bboxes.elements):
                issues = []
                
                if elem.left < 0:
                    issues.append("left < 0")
                if elem.top < 0:
                    issues.append("top < 0")
                if elem.left + elem.width > page_width:
                    issues.append(f"right > {page_width}")
                if elem.top + elem.height > page_height:
                    issues.append(f"bottom > {page_height}")
                
                if issues:
                    out_of_bounds.append({
                        "slide_id": slide_bboxes.slide_id,
                        "element_index": i,
                        "element_type": elem.element_type,
                        "issues": issues,
                        "bbox": elem.to_dict()
                    })
        
        return out_of_bounds
    
    def save_to_json(self, output_path: str):
        """
        保存bounding boxes到JSON文件
        
        Args:
            output_path: 输出文件路径
        """
        data = {
            "pptx_file": str(self.pptx_path),
            "slide_count": len(self.all_bboxes),
            "slides": [s.to_dict() for s in self.all_bboxes],
            "analysis": {
                "overlaps": self.find_overlaps(),
                "out_of_bounds": self.check_boundaries()
            }
        }
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        
        print(f"  ✅ BBox JSON已保存: {output_path}")
    
    def to_dict(self) -> Dict[str, Any]:
        """
        转换为字典格式
        
        Returns:
            包含所有信息的字典
        """
        return {
            "slides": [s.to_dict() for s in self.all_bboxes],
            "overlaps": self.find_overlaps(),
            "out_of_bounds": self.check_boundaries()
        }


def extract_bboxes(pptx_path: str) -> List[SlideBBoxes]:
    """
    便捷函数：提取PPTX的bounding boxes
    
    Args:
        pptx_path: PPTX文件路径
    
    Returns:
        每页的bounding boxes列表
    """
    extractor = BBoxExtractor(pptx_path)
    return extractor.extract()


