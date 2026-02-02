#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
测试脚本：验证文本高度修复
重新提取bbox并对比原始vs修复后的高度
"""

import sys
import json
from pathlib import Path

# 添加项目根目录到 Python 路径
project_root = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(project_root))

from modules.videopipeline.stage4_bbox_extractor import BBoxExtractor
from modules.videopipeline.models import SlideStructure
from pptx import Presentation


def test_text_height_fix():
    """测试文本高度修复"""
    print("=" * 80)
    print("🧪 测试文本高度修复")
    print("=" * 80)
    print()
    
    base_dir = Path(project_root) / "data" / "pipeline_outputs" / "test_ppt_only"
    slides_data_path = base_dir / "slides_data.json"
    pptx_path = base_dir / "pptx" / "presentation.pptx"
    
    # 1. 加载slides_data
    print("📋 加载slides_data...")
    with open(slides_data_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    slides = []
    for slide_dict in data.get('slides', []):
        slide = SlideStructure.from_dict(slide_dict)
        slides.append(slide)
    
    print(f"   ✅ 加载完成，共 {len(slides)} 页")
    print()
    
    # 2. 重新提取bbox
    print("🔍 重新提取BBox（使用文本高度修复）...")
    extractor = BBoxExtractor(str(pptx_path), slides)
    new_bboxes = extractor.extract()  # 正确的方法名
    print(f"   ✅ 提取完成")
    print()
    
    # 3. 对比分析
    print("📊 文本高度对比分析:")
    print()
    
    prs = Presentation(str(pptx_path))
    
    for slide_idx, slide_bboxes in enumerate(new_bboxes):
        if not slide_bboxes.elements:
            continue
        
        print(f"📄 Slide {slide_idx}:")
        
        slide = prs.slides[slide_idx]
        shapes = [s for s in slide.shapes if not s.is_placeholder]
        
        for bbox_idx, bbox in enumerate(slide_bboxes.elements):
            if bbox.element_type != "text":
                continue
            
            # 获取对应的shape
            shape_idx = bbox.shape_index if bbox.shape_index is not None else bbox_idx
            if shape_idx >= len(shapes):
                continue
            
            shape = shapes[shape_idx]
            shape_height = shape.height / 914400
            bbox_height = bbox.height
            
            height_diff = bbox_height - shape_height
            
            print(f"   元素 {bbox_idx} (shape {shape_idx}):")
            print(f"      内容: {bbox.content[:50]}...")
            print(f"      Shape高度: {shape_height:.2f} 英寸")
            print(f"      BBox高度:  {bbox_height:.2f} 英寸")
            
            if abs(height_diff) > 0.05:
                if height_diff > 0:
                    print(f"      ✅ 增加了 {height_diff:.2f} 英寸 (考虑文本实际渲染)")
                else:
                    print(f"      ⚠️ 减少了 {abs(height_diff):.2f} 英寸 (异常)")
            else:
                print(f"      ℹ️  高度未调整 (差异<0.05英寸)")
        
        print()
    
    # 4. 保存新的bboxes
    print("💾 保存新的bboxes...")
    output_path = base_dir / "bboxes_fixed.json"
    
    bboxes_dict = {
        "pptx_file": str(pptx_path),
        "slide_count": len(new_bboxes),
        "slides": [sb.to_dict() for sb in new_bboxes]
    }
    
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(bboxes_dict, f, ensure_ascii=False, indent=2)
    
    print(f"   ✅ 已保存到: {output_path}")
    print()
    
    print("=" * 80)
    print("✅ 测试完成")
    print("=" * 80)
    print()
    print("💡 建议:")
    print("   1. 检查上面的高度对比，确认文本元素的高度是否合理增加")
    print("   2. 使用 bboxes_fixed.json 重新运行VLM优化")
    print("   3. 对比优化前后的重叠情况")
    print()


if __name__ == "__main__":
    test_text_height_fix()

