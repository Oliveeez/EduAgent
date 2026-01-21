"""
Step 3 & 4: PPT创建和编辑模块
支持多轮交互修改
"""

import json
import logging
from typing import Dict, List
from datetime import datetime
from pathlib import Path
import asyncio

from utils.ppt_generator import PPTGenerator
from utils.math_renderer import MathRenderer
from config.config import PPT_OUTPUT_DIR

logger = logging.getLogger(__name__)


class PPTCreator:
    """PPT创建器 - Step 3 & 4"""
    
    def __init__(self, template_path: str = None, llm_client=None):
        """
        初始化PPT创建器
        
        Args:
            template_path: PPT模板路径
            llm_client: LLM客户端
        """
        self.template_path = template_path
        self.llm_client = llm_client
        self.math_renderer = MathRenderer()
        self.current_ppt = None
        self.modification_history = []
        
    async def create_initial_ppt(self, outline: Dict, script: Dict, 
                                latex_data: Dict = None,
                                output_path: str = None) -> Dict:
        """
        创建初版PPT（Step 3）
        
        Args:
            outline: 大纲数据
            script: 讲稿数据
            latex_data: LaTeX数据（包含公式等）
            output_path: 输出路径
            
        Returns:
            创建结果
        """
        logger.info("=" * 60)
        logger.info("🚀 Step 3: 开始创建初版PPT")
        logger.info("=" * 60)
        
        try:
            # 1. 渲染数学公式（如果有）
            math_images = {}
            if latex_data and latex_data.get('equations'):
                logger.info("🔢 Step 3.1: 渲染数学公式...")
                equations = latex_data['equations']
                math_images = self.math_renderer.batch_render(equations)
            
            # 2. 初始化PPT生成器
            logger.info("📑 Step 3.2: 初始化PPT生成器...")
            generator = PPTGenerator(self.template_path)
            
            # 3. 生成PPT
            logger.info("🎨 Step 3.3: 生成PPT内容...")
            self.current_ppt = generator.create_from_outline(outline, script, math_images)
            
            # 4. 保存PPT
            if not output_path:
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                output_path = str(PPT_OUTPUT_DIR / f"presentation_{timestamp}.pptx")
            
            # 确保输出路径是绝对路径
            output_path = str(Path(output_path).resolve())
            Path(output_path).parent.mkdir(exist_ok=True, parents=True)
            
            logger.info("💾 Step 3.4: 保存PPT...")
            saved_path = generator.save(output_path)
            
            result = {
                'success': True,
                'output_path': saved_path,
                'metadata': {
                    'created_at': datetime.now().isoformat(),
                    'template': self.template_path,
                    'total_slides': len(self.current_ppt.slides) if self.current_ppt else 0,
                    'math_images_count': len(math_images)
                }
            }
            
            logger.info("=" * 60)
            logger.info("✅ Step 3: 初版PPT创建完成")
            logger.info(f"   📄 幻灯片数: {result['metadata']['total_slides']}")
            logger.info(f"   💾 保存路径: {saved_path}")
            logger.info("=" * 60)
            
            return result
            
        except Exception as e:
            logger.error(f"❌ 创建PPT失败: {e}")
            import traceback
            traceback.print_exc()
            return {
                'success': False,
                'error': str(e)
            }
    
    async def modify_ppt(self, ppt_path: str, modifications: Dict) -> Dict:
        """
        修改PPT（Step 4）
        
        Args:
            ppt_path: PPT文件路径
            modifications: 修改指令
                - slide_number: 幻灯片编号
                - modification_type: 修改类型 (edit_text, add_image, adjust_layout, etc.)
                - details: 具体修改内容
                
        Returns:
            修改结果
        """
        logger.info("=" * 60)
        logger.info("🚀 Step 4: 开始修改PPT")
        logger.info("=" * 60)
        
        try:
            # 1. 解析修改指令
            logger.info("📝 Step 4.1: 解析修改指令...")
            parsed_modifications = await self._parse_modifications(modifications)
            
            # 2. 应用修改
            logger.info("🔧 Step 4.2: 应用修改...")
            generator = PPTGenerator(ppt_path)
            
            for mod in parsed_modifications:
                self._apply_single_modification(generator, mod)
            
            # 3. 保存修改后的PPT
            logger.info("💾 Step 4.3: 保存修改...")
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            # 确保路径是绝对路径
            ppt_path_abs = str(Path(ppt_path).resolve())
            output_path = ppt_path_abs.replace('.pptx', f'_modified_{timestamp}.pptx')
            saved_path = generator.save(output_path)
            
            # 4. 记录修改历史
            self.modification_history.append({
                'timestamp': datetime.now().isoformat(),
                'modifications': modifications,
                'output_path': saved_path
            })
            
            result = {
                'success': True,
                'output_path': saved_path,
                'modifications_applied': len(parsed_modifications),
                'metadata': {
                    'modified_at': datetime.now().isoformat()
                }
            }
            
            logger.info("=" * 60)
            logger.info("✅ Step 4: PPT修改完成")
            logger.info(f"   🔧 应用了 {len(parsed_modifications)} 项修改")
            logger.info(f"   💾 保存路径: {saved_path}")
            logger.info("=" * 60)
            
            return result
            
        except Exception as e:
            logger.error(f"❌ 修改PPT失败: {e}")
            import traceback
            traceback.print_exc()
            return {
                'success': False,
                'error': str(e)
            }
    
    async def _parse_modifications(self, modifications: Dict) -> List[Dict]:
        """
        解析用户的修改指令（支持自然语言）
        
        如果提供了LLM，可以理解自然语言指令
        否则，直接使用结构化指令
        """
        if isinstance(modifications, list):
            # 已经是结构化指令
            return modifications
        
        # 如果是自然语言描述，使用LLM解析
        if self.llm_client and 'description' in modifications:
            return await self._parse_with_llm(modifications['description'])
        
        # 默认：将单个修改转换为列表
        return [modifications]
    
    async def _parse_with_llm(self, description: str) -> List[Dict]:
        """使用LLM解析自然语言修改指令"""
        prompt = f"""用户想要修改PPT，请将以下自然语言描述转换为结构化的修改指令。

用户描述：
{description}

请返回JSON格式的修改指令列表，格式如下：
[
  {{
    "slide_number": 幻灯片编号(从1开始),
    "modification_type": "修改类型(edit_text/add_image/adjust_layout/split_slide)",
    "details": {{
      "target": "修改目标(title/content/etc)",
      "action": "具体操作",
      "value": "新内容或参数"
    }}
  }}
]

修改类型说明：
- edit_text: 编辑文本内容
- add_image: 添加图片
- adjust_layout: 调整布局
- split_slide: 拆分幻灯片

只返回JSON，不要包含其他说明。"""

        try:
            loop = asyncio.get_event_loop()
            response = await loop.run_in_executor(None, self.llm_client, prompt)
            
            # 清理和解析
            response = response.replace("<br>", "\n").strip()
            import re
            response = re.sub(r'^```json\s*', '', response)
            response = re.sub(r'^```\s*', '', response)
            response = re.sub(r'\s*```$', '', response)
            
            modifications = json.loads(response)
            logger.info(f"✅ LLM解析了 {len(modifications)} 项修改指令")
            return modifications
            
        except Exception as e:
            logger.error(f"LLM解析失败: {e}")
            # 返回默认修改
            return [{
                "slide_number": 1,
                "modification_type": "edit_text",
                "details": {
                    "target": "content",
                    "action": "update",
                    "value": description
                }
            }]
    
    def _apply_single_modification(self, generator: PPTGenerator, modification: Dict):
        """应用单个修改"""
        slide_num = modification.get('slide_number', 1) - 1  # 转换为0-based索引
        mod_type = modification.get('modification_type')
        details = modification.get('details', {})
        
        try:
            prs = generator.prs
            
            if slide_num >= len(prs.slides):
                logger.warning(f"⚠️  幻灯片 {slide_num + 1} 不存在，跳过修改")
                return
            
            slide = prs.slides[slide_num]
            
            if mod_type == 'edit_text':
                self._modify_text(slide, details)
            elif mod_type == 'add_image':
                self._add_image(slide, details)
            elif mod_type == 'adjust_layout':
                self._adjust_layout(slide, details)
            elif mod_type == 'split_slide':
                self._split_slide(generator, slide_num, details)
            else:
                logger.warning(f"⚠️  未知的修改类型: {mod_type}")
            
            logger.info(f"  ✓ 应用修改: 幻灯片{slide_num + 1} - {mod_type}")
            
        except Exception as e:
            logger.error(f"应用修改失败: {e}")
    
    def _modify_text(self, slide, details: Dict):
        """修改文本内容"""
        target = details.get('target', 'content')
        value = details.get('value', '')
        
        if target == 'title' and slide.shapes.title:
            slide.shapes.title.text = value
        elif target == 'content':
            # 修改第一个内容框
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                    shape.text_frame.text = value
                    break
    
    def _add_image(self, slide, details: Dict):
        """添加图片"""
        image_path = details.get('image_path')
        position = details.get('position', {'left': 1, 'top': 1})
        size = details.get('size', {'width': 3})
        
        if not image_path or not Path(image_path).exists():
            logger.warning(f"⚠️  图片不存在: {image_path}")
            return
        
        from pptx.util import Inches
        slide.shapes.add_picture(
            image_path,
            Inches(position['left']),
            Inches(position['top']),
            width=Inches(size['width'])
        )
    
    def _adjust_layout(self, slide, details: Dict):
        """调整布局（简化实现）"""
        # 这里可以实现更复杂的布局调整逻辑
        logger.info(f"调整布局: {details}")
    
    def _split_slide(self, generator: PPTGenerator, slide_num: int, details: Dict):
        """拆分幻灯片"""
        # 这是一个复杂的操作，简化实现
        logger.info(f"拆分幻灯片 {slide_num + 1}")
        # 可以添加一张新幻灯片，并将内容分配到两张幻灯片上
    
    def get_modification_history(self) -> List[Dict]:
        """获取修改历史"""
        return self.modification_history