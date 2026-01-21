"""
PPT生成工具
基于python-pptx库，使用模板生成PPT
"""

import json
import logging
from pathlib import Path
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import shutil

logger = logging.getLogger(__name__)


class PPTGenerator:
    """PPT生成器（基于模板）"""
    
    def __init__(self, template_path: str = None):
        """
        初始化PPT生成器
        
        Args:
            template_path: PPT模板路径
        """
        self.template_path = template_path
        self.prs = None
        
        if template_path and Path(template_path).exists():
            self.prs = Presentation(template_path)
            logger.info(f"📑 加载PPT模板: {template_path}")
            logger.info(f"   模板包含 {len(self.prs.slides)} 张幻灯片")
        else:
            self.prs = Presentation()
            logger.info("📑 创建空白PPT")
    
    def create_from_outline(self, outline: Dict, script: Dict, 
                           math_images: Dict = None) -> str:
        """
        根据大纲和讲稿创建PPT
        
        Args:
            outline: 大纲数据
            script: 讲稿数据
            math_images: 数学公式图片映射
            
        Returns:
            生成的PPT文件路径
        """
        logger.info("🎨 开始生成PPT...")
        
        # 清空现有幻灯片（如果需要）
        if len(self.prs.slides) > 0 and not self.template_path:
            # 对于空白PPT，清空所有幻灯片
            for _ in range(len(self.prs.slides)):
                rId = self.prs.slides._sldIdLst[0].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[0]
        
        # 生成标题页
        self._create_title_slide(outline.get('title', '课程演示'))
        
        # 生成内容页
        sections = outline.get('sections', [])
        for idx, section in enumerate(sections):
            section_script = script.get('sections', [])[idx] if idx < len(script.get('sections', [])) else {}
            self._create_section_slides(section, section_script, math_images)
        
        logger.info(f"✅ PPT生成完成，共 {len(self.prs.slides)} 张幻灯片")
        return self.prs
    
    def _create_title_slide(self, title: str, subtitle: str = None):
        """创建标题页"""
        # 使用第一个布局（通常是标题布局）
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        
        # 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # 设置副标题
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        
        logger.info(f"  ✓ 创建标题页: {title}")
    
    def _create_section_slides(self, section: Dict, script: Dict, math_images: Dict = None):
        """创建章节幻灯片"""
        section_title = section.get('title', '未命名章节')
        points = section.get('points', [])
        
        # 如果章节有多个知识点，可能需要多张幻灯片
        if points:
            # 每个知识点一张幻灯片（或根据内容量决定）
            for point_idx, point in enumerate(points):
                self._create_content_slide(
                    section_title,
                    point,
                    script.get('points', [])[point_idx] if point_idx < len(script.get('points', [])) else {},
                    math_images
                )
        else:
            # 章节概述页
            self._create_content_slide(section_title, section, script, math_images)
    
    def _create_content_slide(self, title: str, content: Dict, 
                             script: Dict = None, math_images: Dict = None):
        """创建内容页"""
        # 使用内容布局（通常是第二个布局）
        if len(self.prs.slide_layouts) > 1:
            content_layout = self.prs.slide_layouts[1]
        else:
            content_layout = self.prs.slide_layouts[0]
        
        slide = self.prs.slides.add_slide(content_layout)
        
        # 设置标题
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # 添加内容
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            
            # 添加知识点内容
            point_content = content.get('content', '')
            examples = content.get('examples', [])
            
            # 主要内容
            if point_content:
                p = text_frame.paragraphs[0]
                p.text = point_content
                p.level = 0
            
            # 示例
            for example in examples:
                p = text_frame.add_paragraph()
                p.text = f"示例: {example}"
                p.level = 1
            
            # 添加公式（如果有）
            equations = content.get('equations', [])
            if equations and math_images:
                for eq_id in equations:
                    if eq_id in math_images and math_images[eq_id]:
                        # 添加公式图片
                        try:
                            left = Inches(1)
                            top = Inches(4)
                            slide.shapes.add_picture(
                                math_images[eq_id],
                                left, top,
                                height=Inches(0.5)
                            )
                        except Exception as e:
                            logger.warning(f"添加公式图片失败: {e}")
        
        # 添加备注（讲稿）
        if script and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            script_text = script.get('text', '')
            if script_text:
                text_frame.text = script_text
        
        logger.info(f"  ✓ 创建内容页: {title}")
    
    def add_custom_slide(self, layout_idx: int = 1) -> object:
        """
        添加自定义幻灯片
        
        Args:
            layout_idx: 布局索引
            
        Returns:
            Slide对象
        """
        layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(layout)
        return slide
    
    def add_text_box(self, slide, text: str, left: float, top: float, 
                    width: float, height: float, font_size: int = 18, bold: bool = False):
        """
        添加文本框
        
        Args:
            slide: 幻灯片对象
            text: 文本内容
            left, top, width, height: 位置和大小（inches）
            font_size: 字体大小
            bold: 是否加粗
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
    
    def add_image(self, slide, image_path: str, left: float, top: float, 
                 width: float = None, height: float = None):
        """
        添加图片
        
        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            left, top: 位置（inches）
            width, height: 大小（inches），可选
        """
        try:
            if width and height:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left), Inches(top),
                    Inches(width), Inches(height)
                )
            elif width:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left), Inches(top),
                    width=Inches(width)
                )
            elif height:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left), Inches(top),
                    height=Inches(height)
                )
            else:
                slide.shapes.add_picture(
                    image_path,
                    Inches(left), Inches(top)
                )
            logger.info(f"  ✓ 添加图片: {image_path}")
        except Exception as e:
            logger.error(f"添加图片失败: {e}")
    
    def save(self, output_path: str) -> str:
        """
        保存PPT
        
        Args:
            output_path: 输出路径
            
        Returns:
            保存的文件路径
        """
        try:
            self.prs.save(output_path)
            logger.info(f"💾 PPT已保存: {output_path}")
            return output_path
        except Exception as e:
            logger.error(f"保存PPT失败: {e}")
            raise


class PPTTemplateManager:
    """PPT模板管理器"""
    
    def __init__(self, templates_dir: str):
        self.templates_dir = Path(templates_dir)
        self.templates_dir.mkdir(exist_ok=True, parents=True)
    
    def list_templates(self) -> List[Dict]:
        """列出所有可用模板"""
        templates = []
        
        for template_file in self.templates_dir.glob("*.pptx"):
            templates.append({
                'name': template_file.stem,
                'path': str(template_file),
                'size': template_file.stat().st_size
            })
        
        return templates
    
    def get_template_info(self, template_path: str) -> Dict:
        """获取模板信息"""
        try:
            prs = Presentation(template_path)
            
            return {
                'path': template_path,
                'slide_count': len(prs.slides),
                'layout_count': len(prs.slide_layouts),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height,
                'layouts': [
                    {
                        'index': idx,
                        'name': layout.name if hasattr(layout, 'name') else f"Layout {idx}"
                    }
                    for idx, layout in enumerate(prs.slide_layouts)
                ]
            }
        except Exception as e:
            logger.error(f"获取模板信息失败: {e}")
            return {}
    
    def copy_template(self, source_path: str, target_name: str = None) -> str:
        """复制模板到模板目录"""
        source = Path(source_path)
        
        if not source.exists():
            raise FileNotFoundError(f"模板文件不存在: {source_path}")
        
        if not target_name:
            target_name = source.name
        
        target_path = self.templates_dir / target_name
        shutil.copy2(source, target_path)
        
        logger.info(f"✅ 模板已复制: {target_path}")
        return str(target_path)